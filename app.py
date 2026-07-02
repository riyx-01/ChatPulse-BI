import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import re
from datetime import datetime, timedelta
from textblob import TextBlob
import io
import random
from fpdf import FPDF
import base64
import zipfile
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import os
import shutil

# Safe import for python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    pptx_installed = True
except ImportError:
    pptx_installed = False

# --- Copy Logo to Workspace ---
logo_source = r"C:\Users\riyat\.gemini\antigravity-ide\brain\5db50240-d5fc-4414-83cd-416f608a2b3b\media__1782980339372.png"
logo_dest = os.path.join(os.path.dirname(__file__), "logo.png")
if os.path.exists(logo_source) and not os.path.exists(logo_dest):
    try:
        shutil.copy(logo_source, logo_dest)
    except:
        pass

st.set_page_config(page_title="ChatPulse | Premium WhatsApp BI", page_icon="logo.png" if os.path.exists(logo_dest) else "📈", layout="wide", initial_sidebar_state="expanded")

# --- Keep Everything in Day Light Only ---
theme_mode = "Day Mode"
bg_gradient = "linear-gradient(135deg, #f8fafc 0%, #e2e8f0 100%)"
card_bg = "rgba(27, 78, 245, 0.02)"
card_border = "rgba(27, 78, 245, 0.15)"
card_hover_bg = "rgba(27, 78, 245, 0.05)"
text_color = "#0f172a"
muted_text = "#475569"
plotly_template = "plotly"
plt_style = "default"
nav_bg = "rgba(255, 255, 255, 0.8)"

# --- Custom CSS for Premium UI & Colors ---
st.markdown(f"""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Plus+Jakarta+Sans:wght@300;400;500;600;700;800&display=swap');

    html, body, [class*="css"] {{
        font-family: 'Plus Jakarta Sans', sans-serif;
    }}
    
    .stApp {{
        background: {bg_gradient};
        color: {text_color};
    }}
    
    .navbar {{
        display: flex;
        align-items: center;
        gap: 20px;
        background: {nav_bg};
        backdrop-filter: blur(12px);
        border-bottom: 1px solid rgba(56, 116, 255, 0.2);
        padding: 15px 40px;
        border-radius: 12px;
        margin-bottom: 30px;
    }}
    .navbar-brand {{
        font-size: 1.8rem;
        font-weight: 800;
        background: linear-gradient(90deg, #1B4EF5, #5996FF, #F4CEFF);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        letter-spacing: -0.5px;
    }}
    
    div[data-testid="metric-container"] {{
        background: {card_bg};
        backdrop-filter: blur(12px);
        -webkit-backdrop-filter: blur(12px);
        border: 1px solid {card_border};
        padding: 24px;
        border-radius: 16px;
        box-shadow: 0 10px 30px rgba(0, 0, 0, 0.05);
        transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
    }}
    div[data-testid="metric-container"]:hover {{
        transform: translateY(-5px);
        border-color: rgba(89, 150, 255, 0.4);
        background: {card_hover_bg};
        box-shadow: 0 15px 35px rgba(27, 78, 245, 0.1);
    }}
    div[data-testid="metric-container"] label {{
        color: {muted_text};
        font-weight: 600;
        font-size: 0.9rem;
        letter-spacing: 0.5px;
        text-transform: uppercase;
    }}
    div[data-testid="metric-container"] div {{
        color: {text_color};
        font-weight: 800;
        font-size: 2.2rem;
    }}
    
    .stButton>button {{
        background: linear-gradient(135deg, #1B4EF5 0%, #3874FF 100%) !important;
        color: #ffffff !important;
        border: none !important;
        border-radius: 8px !important;
        padding: 0.6rem 1.5rem !important;
        font-weight: 600 !important;
        box-shadow: 0 4px 15px rgba(27, 78, 245, 0.2);
        transition: all 0.2s ease;
    }}
    .stButton>button:hover {{
        transform: translateY(-2px);
        box-shadow: 0 6px 20px rgba(56, 116, 255, 0.3);
    }}
    
    .stTabs [data-baseweb="tab-list"] {{
        gap: 15px;
        background-color: rgba(0, 0, 0, 0.02);
        padding: 8px;
        border-radius: 12px;
        border: 1px solid rgba(0, 0, 0, 0.05);
    }}
    .stTabs [data-baseweb="tab"] {{
        background-color: transparent;
        border-radius: 8px;
        color: {muted_text};
        padding: 10px 20px;
        font-weight: 600;
        transition: all 0.3s;
    }}
    .stTabs [aria-selected="true"] {{
        background: #1B4EF5 !important;
        color: white !important;
        box-shadow: 0 4px 12px rgba(27, 78, 245, 0.2);
    }}
    
    .insight-card {{
        background: rgba(27, 78, 245, 0.05);
        border-left: 4px solid #1B4EF5;
        padding: 15px;
        border-radius: 4px;
        margin-top: 15px;
        font-size: 0.95rem;
        color: #1e293b;
    }}
</style>
""", unsafe_allow_html=True)

# --- Global Navigation Bar with Logo ---
logo_html = ""
if os.path.exists(logo_dest):
    with open(logo_dest, "rb") as image_file:
        encoded_logo = base64.b64encode(image_file.read()).decode()
    logo_html = f'<img src="data:image/png;base64,{encoded_logo}" style="height: 50px; border-radius: 8px;">'

st.markdown(f"""
<div class="navbar">
    {logo_html}
    <div class="navbar-brand">ChatPulse BI</div>
    <div style="color: {muted_text}; font-size: 0.95rem; font-weight: 500;">
        Professional WhatsApp Analytics & Intelligence Dashboard
    </div>
</div>
""", unsafe_allow_html=True)

# --- PDF Unicode/Emoji Sanitizer ---
def sanitize_pdf_text(text):
    """Filters string to fit inside PDF standard Latin-1 encoding, replacing emojis with equivalent placeholders."""
    cleaned = []
    for char in str(text):
        if ord(char) < 256:
            cleaned.append(char)
        else:
            cleaned.append("?")
    return "".join(cleaned)

# --- Matplotlib Plotting Helpers for Reports ---
def plot_messages_per_day(ax, df):
    ax.plot(df['full_date'], df['count'], color='#1B4EF5', linewidth=2)
    ax.fill_between(df['full_date'], df['count'], color='#1B4EF5', alpha=0.3)
    ax.set_ylabel("Messages")
    ax.set_title("Messages per Day Over Time")

def plot_messages_per_hour(ax, df):
    ax.bar(df['hour_of_day'], df['count'], color='#3874FF')
    ax.set_xlabel("Hour of Day (0-23)")
    ax.set_ylabel("Messages")
    ax.set_title("Messages by Hour of Day")

def plot_activity_heatmap(ax, df):
    im = ax.imshow(df.values, cmap="RdYlGn_r", aspect='auto')
    ax.set_xticks(range(len(df.columns)))
    ax.set_xticklabels(df.columns)
    ax.set_yticks(range(len(df.index)))
    ax.set_yticklabels(df.index)
    plt.colorbar(im, ax=ax)
    ax.set_title("Activity Heatmap (Day vs Hour)")

def plot_user_volume(ax, df):
    names = [sanitize_pdf_text(name) for name in df['user_name']]
    ax.barh(names, df['count'], color='#5996FF')
    ax.set_xlabel("Messages")
    ax.set_title("Message Count per User")

def plot_convo_share(ax, df, palette):
    names = [sanitize_pdf_text(name) for name in df['user_name']]
    ax.pie(df['count'], labels=names, colors=palette, autopct='%1.1f%%')
    ax.set_title("Conversation Share")

def plot_avg_response_time(ax, df):
    names = [sanitize_pdf_text(name) for name in df['user_name']]
    ax.bar(names, df['avg_resp'], color='#1B4EF5')
    ax.set_ylabel("Minutes")
    ax.set_title("Average Response Time")

def plot_median_response_time(ax, df):
    names = [sanitize_pdf_text(name) for name in df['user_name']]
    ax.bar(names, df['median_resp'], color='#3874FF')
    ax.set_ylabel("Minutes")
    ax.set_title("Median Response Time")

def plot_rolling_sentiment(ax, df):
    ax.plot(df['full_date'], df['rolling_7d'], color='#5996FF', linewidth=2)
    ax.set_ylabel("Sentiment")
    ax.set_title("Rolling 7-Day Average Sentiment")

def plot_avg_sentiment(ax, df):
    names = [sanitize_pdf_text(name) for name in df['user_name']]
    ax.bar(names, df['sentiment_score'], color='#F4CEFF')
    ax.set_ylabel("Sentiment Score")
    ax.set_title("Average Sentiment by User")

def plot_starters(ax, df, palette):
    names = [sanitize_pdf_text(name) for name in df['user_name']]
    ax.pie(df['count'], labels=names, colors=palette, autopct='%1.1f%%')
    ax.set_title("Conversation Starters")

def plot_weekend_vs_weekday(ax, df):
    ax.bar(df['day_type'], df['count'], color=['#1B4EF5', '#3874FF'])
    ax.set_ylabel("Messages")
    ax.set_title("Weekend vs Weekday")

# --- Report Creators ---
def generate_pdf_report_with_fallback(kpi_dict, data_frames, user_stats, insights):
    """Generates a highly structured and readable PDF report with Matplotlib charts and data tables."""
    pdf = FPDF()
    pdf.add_page()
    
    # Title Block
    pdf.set_fill_color(27, 78, 245)
    pdf.rect(0, 0, 210, 40, 'F')
    pdf.set_text_color(255, 255, 255)
    pdf.set_font("Arial", 'B', 24)
    pdf.cell(0, 20, sanitize_pdf_text("ChatPulse BI Insights Report"), ln=True, align='C')
    pdf.ln(15)
    
    # Executive Summary Box
    pdf.set_text_color(0, 0, 0)
    pdf.set_font("Arial", 'B', 16)
    pdf.cell(0, 10, sanitize_pdf_text("Executive Summary Statistics"), ln=True)
    pdf.set_font("Arial", '', 12)
    pdf.ln(2)
    
    for kpi, val in kpi_dict.items():
        pdf.cell(0, 8, sanitize_pdf_text(f" - {kpi}: {val}"), ln=True)
    pdf.ln(10)
    
    # Detailed User Engagement Table
    pdf.add_page()
    pdf.set_font("Arial", 'B', 16)
    pdf.cell(0, 10, sanitize_pdf_text("User Engagement & Behavior Statistics"), ln=True)
    pdf.set_font("Arial", '', 10)
    pdf.ln(5)
    
    # Table Headers
    pdf.set_fill_color(27, 78, 245)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(40, 10, "User Name", 1, 0, 'C', True)
    pdf.cell(25, 10, "Messages", 1, 0, 'C', True)
    pdf.cell(25, 10, "Avg Len", 1, 0, 'C', True)
    pdf.cell(30, 10, "Avg Resp (m)", 1, 0, 'C', True)
    pdf.cell(30, 10, "Med Resp (m)", 1, 0, 'C', True)
    pdf.cell(30, 10, "Avg Sentiment", 1, 1, 'C', True)
    
    pdf.set_text_color(0, 0, 0)
    for index, row in user_stats.iterrows():
        pdf.cell(40, 10, sanitize_pdf_text(row['user_name']), 1, 0, 'L')
        pdf.cell(25, 10, f"{row['total_messages']:,}", 1, 0, 'C')
        pdf.cell(25, 10, f"{row['avg_message_length']:.1f}", 1, 0, 'C')
        pdf.cell(30, 10, f"{row['avg_response_time']:.1f}", 1, 0, 'C')
        pdf.cell(30, 10, f"{row['median_response_time']:.1f}", 1, 0, 'C')
        pdf.cell(30, 10, f"{row['avg_sentiment']:.2f}", 1, 1, 'C')
    pdf.ln(15)

    # Plot helper function to add page with plot
    def add_plot_to_pdf(title, plot_func, df_arg, insight_text, extra_palette=None):
        pdf.add_page()
        pdf.set_font("Arial", 'B', 16)
        pdf.cell(0, 10, sanitize_pdf_text(title), ln=True, align='C')
        pdf.ln(5)
        
        fig, ax = plt.subplots(figsize=(8, 4.5))
        plt.style.use('default') 
        
        if extra_palette:
            plot_func(ax, df_arg, extra_palette)
        else:
            plot_func(ax, df_arg)
            
        img_buffer = io.BytesIO()
        fig.savefig(img_buffer, format='png', bbox_inches='tight', dpi=150)
        img_buffer.seek(0)
        plt.close(fig)
        
        pdf.image(img_buffer, x=15, y=30, w=180)
        
        # Add written insight text below chart
        pdf.set_y(175)
        pdf.set_font("Arial", 'B', 12)
        pdf.cell(0, 10, sanitize_pdf_text("Business Intelligence Insight:"), ln=True)
        pdf.set_font("Arial", '', 11)
        pdf.multi_cell(0, 6, sanitize_pdf_text(insight_text))

    palette = ['#1B4EF5', '#3874FF', '#5996FF', '#F4CEFF']
    
    # Append fallback graphs to PDF safely
    add_plot_to_pdf("Messages per Day Over Time", plot_messages_per_day, data_frames["daily_vol"], insights["daily_vol"])
    add_plot_to_pdf("Messages by Hour of Day", plot_messages_per_hour, data_frames["hourly_vol"], insights["hourly_vol"])
    add_plot_to_pdf("Activity Heatmap", plot_activity_heatmap, data_frames["heatmap_pivot"], insights["heatmap"])
    add_plot_to_pdf("Message Count per User", plot_user_volume, data_frames["user_vol"], insights["user_vol"])
    add_plot_to_pdf("Conversation Share", plot_convo_share, data_frames["user_vol"], insights["user_share"], palette)
    add_plot_to_pdf("Average Response Time by User", plot_avg_response_time, data_frames["avg_resp"], insights["response_time"])
    add_plot_to_pdf("Rolling 7-Day Average Sentiment", plot_rolling_sentiment, data_frames["daily_sentiment"], insights["sentiment_trend"])
    add_plot_to_pdf("Average Sentiment by User", plot_avg_sentiment, data_frames["user_sent"], insights["sentiment_compare"])
    add_plot_to_pdf("Conversation Starters", plot_starters, data_frames["starters"], insights["starters"], palette)
    add_plot_to_pdf("Weekend vs Weekday", plot_weekend_vs_weekday, data_frames["weekend_vol"], insights["weekend_vol"])

    return pdf.output(dest='S')


def generate_pptx_report(kpi_dict, data_frames, user_stats, insights):
    """Generates a PowerPoint presentation using python-pptx in Day light theme."""
    if not pptx_installed:
        return None
        
    prs = Presentation()
    blank_layout = prs.slide_layouts[6] 
    
    # Light background for Slide
    def set_light_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 250, 252) # Day background color
        
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(blank_layout)
    set_light_background(slide)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9.0), Inches(2.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "ChatPulse BI Insights"
    p.font.bold = True
    p.font.size = Pt(46)
    p.font.color.rgb = RGBColor(27, 78, 245)
    
    p2 = tf.add_paragraph()
    p2.text = "Comprehensive WhatsApp Analytics & Intelligence Report"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(71, 85, 105)
    p2.space_before = Pt(10)
    
    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(blank_layout)
    set_light_background(slide)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(1.0))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Executive Summary"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(27, 78, 245)
    
    txBox_kpis = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.0))
    tf_kpis = txBox_kpis.text_frame
    for kpi, val in kpi_dict.items():
        p_kpi = tf_kpis.add_paragraph()
        p_kpi.text = sanitize_pdf_text(f"{kpi}: {val}")
        p_kpi.font.size = Pt(20)
        p_kpi.font.color.rgb = RGBColor(15, 23, 42)
        p_kpi.space_after = Pt(16)

    # Slide 3: Detailed User Statistics Table
    slide = prs.slides.add_slide(blank_layout)
    set_light_background(slide)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(1.0))
    p = txBox.text_frame.paragraphs[0]
    p.text = "User Behavior & Engagement"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(27, 78, 245)
    
    rows = len(user_stats) + 1
    cols = 6
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(4.5)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    headers = ["User Name", "Messages", "Avg Len", "Avg Resp (m)", "Med Resp (m)", "Avg Sent"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(27, 78, 245)
        
    for row_idx, row in user_stats.iterrows():
        table.cell(row_idx+1, 0).text = sanitize_pdf_text(row['user_name'])
        table.cell(row_idx+1, 1).text = f"{row['total_messages']:,}"
        table.cell(row_idx+1, 2).text = f"{row['avg_message_length']:.1f}"
        table.cell(row_idx+1, 3).text = f"{row['avg_response_time']:.1f}"
        table.cell(row_idx+1, 4).text = f"{row['median_response_time']:.1f}"
        table.cell(row_idx+1, 5).text = f"{row['avg_sentiment']:.2f}"

    # Helper function to add slides with chart images and insights
    def add_chart_slide(title, plot_func, df_arg, insight_text, extra_palette=None):
        slide = prs.slides.add_slide(blank_layout)
        set_light_background(slide)
        
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.8))
        p = txBox.text_frame.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(27, 78, 245)
        
        fig, ax = plt.subplots(figsize=(6, 3.8))
        plt.style.use('default') 
        
        if extra_palette:
            plot_func(ax, df_arg, extra_palette)
        else:
            plot_func(ax, df_arg)
            
        img_buffer = io.BytesIO()
        fig.savefig(img_buffer, format='png', bbox_inches='tight', dpi=150)
        img_buffer.seek(0)
        plt.close(fig)
        
        slide.shapes.add_picture(img_buffer, Inches(0.5), Inches(1.1), Inches(5.8), Inches(3.8))
        
        # Add Insight Text box on the right of the slide
        txBox_insight = slide.shapes.add_textbox(Inches(6.5), Inches(1.1), Inches(3.0), Inches(3.8))
        tf_ins = txBox_insight.text_frame
        tf_ins.word_wrap = True
        
        p_ins_title = tf_ins.paragraphs[0]
        p_ins_title.text = "BI Insight:"
        p_ins_title.font.bold = True
        p_ins_title.font.size = Pt(16)
        p_ins_title.font.color.rgb = RGBColor(27, 78, 245)
        p_ins_title.space_after = Pt(10)
        
        p_ins_desc = tf_ins.add_paragraph()
        p_ins_desc.text = sanitize_pdf_text(insight_text)
        p_ins_desc.font.size = Pt(14)
        p_ins_desc.font.color.rgb = RGBColor(71, 85, 105)

    palette = ['#1B4EF5', '#3874FF', '#5996FF', '#F4CEFF']
    
    # Add Visualizations Slides
    add_chart_slide("Messages per Day Over Time", plot_messages_per_day, data_frames["daily_vol"], insights["daily_vol"])
    add_chart_slide("Messages by Hour of Day", plot_messages_per_hour, data_frames["hourly_vol"], insights["hourly_vol"])
    add_chart_slide("Activity Heatmap", plot_activity_heatmap, data_frames["heatmap_pivot"], insights["heatmap"])
    add_chart_slide("Message Count per User", plot_user_volume, data_frames["user_vol"], insights["user_vol"])
    add_chart_slide("Conversation Share", plot_convo_share, data_frames["user_vol"], insights["user_share"], palette)
    add_chart_slide("Average Response Time by User", plot_avg_response_time, data_frames["avg_resp"], insights["response_time"])
    add_chart_slide("Rolling 7-Day Average Sentiment", plot_rolling_sentiment, data_frames["daily_sentiment"], insights["sentiment_trend"])
    add_chart_slide("Average Sentiment by User", plot_avg_sentiment, data_frames["user_sent"], insights["sentiment_compare"])
    add_chart_slide("Conversation Starters", plot_starters, data_frames["starters"], insights["starters"], palette)
    add_chart_slide("Weekend vs Weekday", plot_weekend_vs_weekday, data_frames["weekend_vol"], insights["weekend_vol"])
    
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    return ppt_buffer.getvalue()


# --- Zip Handling & Parser ---
def extract_zip(uploaded_zip):
    try:
        with zipfile.ZipFile(uploaded_zip) as z:
            for filename in z.namelist():
                if filename.endswith('.txt'):
                    with z.open(filename) as f:
                        return f.read().decode('utf-8', errors='ignore')
        return None
    except Exception as e:
        st.error(f"Error extracting ZIP: {e}")
        return None

@st.cache_data
def parse_chat(file_contents: str) -> pd.DataFrame:
    """Robust parser for WhatsApp (iOS, Android, Web)."""
    pattern1 = re.compile(r'^(\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4}),?\s(\d{1,2}:\d{2}(?::\d{2})?(?:\s?[APap][mM])?)(?:\s-\s|\s)(.*)$')
    pattern2 = re.compile(r'^\[(\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4}),?\s(\d{1,2}:\d{2}(?::\d{2})?(?:\s?[APap][mM])?)\]\s?(.*)$')
    pattern3 = re.compile(r'^\[(\d{1,2}:\d{2}(?::\d{2})?(?:\s?[APap][mM])?),?\s(\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4})\]\s?(.*)$')
    pattern4 = re.compile(r'^(\d{1,2}:\d{2}(?::\d{2})?(?:\s?[APap][mM])?),?\s(\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4})(?:\s-\s|\s)(.*)$')

    data = []
    lines = file_contents.split('\n')
    
    current_date = None
    current_time = None
    current_sender = None
    current_message = ""
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        match = pattern1.match(line)
        if match:
            date_str, time_str, rest = match.groups()
        else:
            match = pattern2.match(line)
            if match:
                date_str, time_str, rest = match.groups()
            else:
                match = pattern3.match(line)
                if match:
                    time_str, date_str, rest = match.groups()
                else:
                    match = pattern4.match(line)
                    if match:
                        time_str, date_str, rest = match.groups()
        
        if match:
            if current_sender and current_message:
                data.append([current_date, current_time, current_sender, current_message])
                
            if ':' in rest:
                parts = rest.split(':', 1)
                current_date = date_str
                current_time = time_str
                current_sender = parts[0].strip()
                current_message = parts[1].strip()
            else:
                current_sender = None
                current_message = ""
        else:
            if current_sender is not None:
                current_message += " " + line
                
    if current_sender and current_message:
        data.append([current_date, current_time, current_sender, current_message])
        
    df = pd.DataFrame(data, columns=['Date', 'Time', 'Sender', 'Message'])
    
    if not df.empty:
        try:
            df['DateTime'] = pd.to_datetime(df['Date'] + ' ' + df['Time'], dayfirst=True, format='mixed', errors='coerce')
        except:
            df['DateTime'] = pd.to_datetime(df['Date'] + ' ' + df['Time'], errors='coerce')
            
        df = df.dropna(subset=['DateTime'])
        df = df.sort_values('DateTime').reset_index(drop=True)
        
        df = df[df['Message'].str.strip() != '']
        df = df[~df['Sender'].str.contains('changed the subject to', na=False, case=False)]
        df = df[~df['Sender'].str.contains('changed the group description', na=False, case=False)]
        
    return df

@st.cache_data
def build_data_model(df: pd.DataFrame):
    if df is None or df.empty:
        return None, None, None
        
    df['message_id'] = range(1, len(df) + 1)
    df['date_key'] = df['DateTime'].dt.strftime('%Y%m%d').astype(int)
    df['message_length'] = df['Message'].str.len()
    df['word_count'] = df['Message'].apply(lambda x: len(str(x).split()))
    
    def get_sentiment(text):
        try:
            return TextBlob(str(text)).sentiment.polarity
        except:
            return 0.0
            
    df['sentiment_score'] = df['Message'].apply(get_sentiment)
    
    df['prev_sender'] = df['Sender'].shift(1)
    df['prev_time'] = df['DateTime'].shift(1)
    
    df['is_reply'] = (df['Sender'] != df['prev_sender']) & (df['prev_sender'].notna())
    
    df['response_time_minutes'] = 0.0
    mask = df['is_reply']
    df.loc[mask, 'response_time_minutes'] = (df.loc[mask, 'DateTime'] - df.loc[mask, 'prev_time']).dt.total_seconds() / 60.0
    
    df['hour_of_day'] = df['DateTime'].dt.hour
    df['is_media'] = df['Message'].str.contains('<Media omitted>', case=False, na=False)
    
    user_stats = df.groupby('Sender').agg(
        first_seen_date=('DateTime', 'min'),
        total_messages=('message_id', 'count')
    ).reset_index()
    
    user_stats['user_key'] = range(1, len(user_stats) + 1)
    user_stats.rename(columns={'Sender': 'user_name'}, inplace=True)
    dim_user = user_stats[['user_key', 'user_name', 'first_seen_date', 'total_messages']]
    
    df = df.merge(user_stats[['user_name', 'user_key']], left_on='Sender', right_on='user_name', how='left')
    
    dates = df['DateTime'].dt.floor('D').unique()
    dim_date = pd.DataFrame({'full_date': dates})
    dim_date['date_key'] = dim_date['full_date'].dt.strftime('%Y%m%d').astype(int)
    dim_date['day_name'] = dim_date['full_date'].dt.day_name()
    dim_date['is_weekend'] = dim_date['full_date'].dt.dayofweek >= 5
    dim_date['week_number'] = dim_date['full_date'].dt.isocalendar().week
    dim_date['month_name'] = dim_date['full_date'].dt.month_name()
    
    fact_message = df[['message_id', 'date_key', 'user_key', 'DateTime', 'message_length', 'word_count', 
                       'sentiment_score', 'is_reply', 'response_time_minutes', 'hour_of_day', 'Message', 'is_media']]
                       
    return dim_date, dim_user, fact_message

def generate_sample_chat():
    users = ["Alice", "Bob", "Charlie", "Diana"]
    messages = [
        "Hey everyone!", "How are we doing?", "I'm good, thanks!", "Did you see the latest update?",
        "Yes, it looks amazing.", "I have some feedback.", "Let's schedule a call.", "Sounds good to me.",
        "Can we push it to tomorrow?", "Sure, no problem.", "<Media omitted>", "Haha nice!",
        "I'm not sure about this.", "This is brilliant!", "I am very disappointed with the results.",
        "Okay.", "Got it.", "Who is joining the meeting?", "I will be 5 mins late."
    ]
    
    lines = []
    current_time = datetime(2023, 1, 1, 9, 0)
    
    for i in range(1000):
        gap = random.expovariate(1/30) 
        current_time += timedelta(minutes=gap)
        sender = random.choices(users, weights=[0.4, 0.3, 0.2, 0.1])[0]
        msg = random.choice(messages)
        time_str = current_time.strftime("%d/%m/%y, %H:%M")
        lines.append(f"{time_str} - {sender}: {msg}")
        
    return "\n".join(lines)


# --- Sidebar ---
with st.sidebar:
    st.markdown("### Data Source")
    uploaded_file = st.file_uploader("Upload WhatsApp Export (.txt or .zip)", type=['txt', 'zip'])
    
    st.markdown("---")
    use_sample = st.button("Use Sample Data Demo", use_container_width=True)
    
    st.markdown("---")
    
    with st.expander("How to Export Chats"):
        st.markdown("""
        **From Android:**
        1. Open the WhatsApp chat.
        2. Tap the three dots (Menu) > **More** > **Export chat**.
        3. Select **Without media**.
        
        **From iOS (iPhone):**
        1. Open the WhatsApp chat.
        2. Tap the contact/group name at the top.
        3. Scroll down and tap **Export Chat**.
        4. Choose **Without Media**.
        
        **From WhatsApp Web:**
        - WhatsApp Web does not natively export chats. However, you can select the text, copy/paste it into a `.txt` file, or use a browser extension to export the conversation log.
        """)

# Process Data
df_raw = None

if uploaded_file is not None:
    try:
        if uploaded_file.name.endswith('.zip'):
            content = extract_zip(uploaded_file)
            if content is None:
                st.error("Could not find a valid .txt file inside the ZIP.")
            else:
                df_raw = parse_chat(content)
        else:
            content = uploaded_file.getvalue().decode("utf-8")
            df_raw = parse_chat(content)
    except Exception as e:
        st.error(f"Error reading file: {str(e)}")
elif use_sample:
    content = generate_sample_chat()
    df_raw = parse_chat(content)

custom_palette = ['#1B4EF5', '#3874FF', '#5996FF', '#F4CEFF']

if df_raw is not None and not df_raw.empty:
    with st.spinner("Processing Data Model & Running Analytics..."):
        dim_date, dim_user, fact_message = build_data_model(df_raw)
        
    if fact_message.empty:
        st.error("Could not parse messages. Please ensure the file is a valid WhatsApp export.")
        st.stop()
        
    # --- KPI CARDS ---
    col1, col2, col3, col4 = st.columns(4)
    
    total_messages = len(fact_message)
    unique_users = len(dim_user)
    num_days = len(dim_date)
    avg_msg_per_day = round(total_messages / num_days) if num_days > 0 else total_messages
    
    user_counts = fact_message.groupby('user_key').size()
    if not user_counts.empty:
        top_user_key = user_counts.idxmax()
        most_active_user = dim_user[dim_user['user_key'] == top_user_key]['user_name'].iloc[0]
    else:
        most_active_user = "N/A"
        
    peak_hour = fact_message['hour_of_day'].mode()
    peak_hour_str = f"{peak_hour.iloc[0]}:00" if not peak_hour.empty else "N/A"
    avg_sentiment = fact_message['sentiment_score'].mean()
    
    col1.metric("Total Messages", f"{total_messages:,}")
    col2.metric("Unique Users", f"{unique_users}")
    col3.metric("Avg Messages / Day", f"{avg_msg_per_day}")
    col4.metric("Most Active User", most_active_user)
    
    st.markdown("<br>", unsafe_allow_html=True)
    col1b, col2b, col3b, col4b = st.columns(4)
    col1b.metric("Peak Activity Hour", peak_hour_str)
    col2b.metric("Avg Sentiment Score", f"{avg_sentiment:.2f}")
    
    data_frames = {}
    
    fact_with_date = fact_message.merge(dim_date, on='date_key')
    fact_with_user = fact_message.merge(dim_user, on='user_key')
    
    # Analytics data preparation
    daily_vol = fact_with_date.groupby('full_date').size().reset_index(name='count')
    data_frames["daily_vol"] = daily_vol
    
    hourly_vol = fact_message.groupby('hour_of_day').size().reset_index(name='count')
    data_frames["hourly_vol"] = hourly_vol
    
    heatmap_data = fact_with_date.groupby(['day_name', 'hour_of_day']).size().reset_index(name='count')
    heatmap_pivot = heatmap_data.pivot(index='day_name', columns='hour_of_day', values='count').fillna(0)
    heatmap_pivot = heatmap_pivot.reindex(['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday'])
    data_frames["heatmap_pivot"] = heatmap_pivot
    
    user_vol = fact_with_user.groupby('user_name').size().reset_index(name='count').sort_values('count', ascending=True)
    data_frames["user_vol"] = user_vol
    
    # User Engagement Stats Table
    user_stats = fact_with_user.groupby('user_name').agg(
        total_messages=('message_id', 'count'),
        avg_message_length=('message_length', 'mean'),
        avg_response_time=('response_time_minutes', lambda x: x[x>0].mean()),
        median_response_time=('response_time_minutes', lambda x: x[x>0].median()),
        avg_sentiment=('sentiment_score', 'mean')
    ).reset_index()
    user_days = fact_with_user.groupby(['user_name', 'date_key']).size().reset_index()
    user_days_count = user_days.groupby('user_name').size().reset_index(name='active_days')
    user_stats = user_stats.merge(user_days_count, on='user_name')
    user_stats['msgs_per_active_day'] = user_stats['total_messages'] / user_stats['active_days']
    user_stats = user_stats.round(2).sort_values('total_messages', ascending=False).reset_index(drop=True)

    # Response times
    valid_responses = fact_with_user[fact_with_user['response_time_minutes'] > 0]
    avg_resp = valid_responses.groupby('user_name')['response_time_minutes'].mean().reset_index(name='avg_resp').sort_values('avg_resp')
    median_resp = valid_responses.groupby('user_name')['response_time_minutes'].median().reset_index(name='median_resp').sort_values('median_resp')
    data_frames["avg_resp"] = avg_resp
    data_frames["median_resp"] = median_resp
    
    # Sentiment
    daily_sentiment = fact_with_date.groupby('full_date')['sentiment_score'].mean().reset_index()
    daily_sentiment['rolling_7d'] = daily_sentiment['sentiment_score'].rolling(window=7, min_periods=1).mean()
    data_frames["daily_sentiment"] = daily_sentiment
    
    user_sent = fact_with_user.groupby('user_name')['sentiment_score'].mean().reset_index().sort_values('sentiment_score')
    data_frames["user_sent"] = user_sent
    
    # Starters
    first_msgs = fact_with_user.loc[fact_with_user.groupby('date_key')['DateTime'].idxmin()]
    starters = first_msgs['user_name'].value_counts().reset_index()
    starters.columns = ['user_name', 'count']
    data_frames["starters"] = starters
    
    # Weekend
    weekend_vol = fact_with_date.groupby('is_weekend').size().reset_index(name='count')
    weekend_vol['day_type'] = weekend_vol['is_weekend'].map({True: 'Weekend', False: 'Weekday'})
    data_frames["weekend_vol"] = weekend_vol
    
    # --- DYNAMIC BUSINESS INTELLIGENCE INSIGHTS ---
    insights = {}
    
    # 1. Daily Volume Insight
    max_day_row = daily_vol.loc[daily_vol['count'].idxmax()]
    insights["daily_vol"] = f"Activity peaked on {max_day_row['full_date'].strftime('%Y-%m-%d')} with {max_day_row['count']} messages sent. The line graph shows communication trends and patterns over time."
    
    # 2. Hourly Volume Insight
    peak_hr = hourly_vol.loc[hourly_vol['count'].idxmax()]['hour_of_day']
    insights["hourly_vol"] = f"Communication is most concentrated at hour {peak_hr}:00, indicating a clear peak time for group interactions during this part of the day."
    
    # 3. Heatmap Insight
    insights["heatmap"] = f"The density distribution shows active interaction slots. Darker areas represent the absolute highest-volume hours in the week."
    
    # 4. User Volume Insight
    if len(user_stats) > 0:
        top_u = user_stats.iloc[0]['user_name']
        top_u_count = user_stats.iloc[0]['total_messages']
        insights["user_vol"] = f"{top_u} is the primary messaging driver with {top_u_count:,} messages. There is a clear divide in engagement levels among active users."
        insights["user_share"] = f"{top_u} controls the largest share of voice in the group. This pie chart highlights who dominates the overall conversation flow."
    else:
        insights["user_vol"] = "No user metrics available."
        insights["user_share"] = "No user conversation share metrics available."
        
    # 5. Response Time Insight
    if len(avg_resp) > 0:
        fast_u = avg_resp.iloc[0]['user_name']
        fast_t = avg_resp.iloc[0]['avg_resp']
        slow_u = avg_resp.iloc[-1]['user_name']
        slow_t = avg_resp.iloc[-1]['avg_resp']
        insights["response_time"] = f"{fast_u} responds the quickest (avg {fast_t:.1f} minutes), while {slow_u} has the longest delay (avg {slow_t:.1f} minutes)."
    else:
        insights["response_time"] = "No response metrics available."
        
    # 6. Sentiment Trend Insight
    insights["sentiment_trend"] = "The rolling sentiment chart highlights the overall group mood progression. Upswings indicate positive interaction periods."
    
    # 7. Sentiment Compare Insight
    if len(user_sent) > 0:
        pos_u = user_sent.iloc[-1]['user_name']
        pos_v = user_sent.iloc[-1]['sentiment_score']
        neg_u = user_sent.iloc[0]['user_name']
        neg_v = user_sent.iloc[0]['sentiment_score']
        insights["sentiment_compare"] = f"{pos_u} exhibits the most positive sentiment index ({pos_v:.2f}), whereas {neg_u} averages the most negative index ({neg_v:.2f})."
    else:
        insights["sentiment_compare"] = "No sentiment metrics available."
        
    # 8. Starters Insight
    if len(starters) > 0:
        top_start = starters.iloc[0]['user_name']
        top_start_count = starters.iloc[0]['count']
        insights["starters"] = f"{top_start} is the primary conversation catalyst, initiating the group chat {top_start_count} times."
    else:
        insights["starters"] = "No starter stats available."
        
    # 9. Weekend Vol Insight
    weekday_c = weekend_vol[weekend_vol['day_type']=='Weekday']['count'].values
    weekend_c = weekend_vol[weekend_vol['day_type']=='Weekend']['count'].values
    wd_val = weekday_c[0] if len(weekday_c)>0 else 1
    we_val = weekend_c[0] if len(weekend_c)>0 else 0
    insights["weekend_vol"] = f"Weekday volume is {wd_val:,} messages vs Weekend volume of {we_val:,} messages. Communication shifts significantly based on the workweek schedule."

    st.markdown("<br>", unsafe_allow_html=True)
    
    # --- APP TABS ---
    tab1, tab2, tab3, tab4 = st.tabs(["Activity Trends", "User Analytics", "Sentiment & Patterns", "Generate Report"])
    
    with tab1:
        c1, c2 = st.columns(2)
        fig_daily = px.line(daily_vol, x='full_date', y='count', title="Messages per Day Over Time", template=plotly_template, color_discrete_sequence=custom_palette)
        c1.plotly_chart(fig_daily, use_container_width=True)
        c1.markdown(f"<div class='insight-card'><b>Analysis Insight:</b> {insights['daily_vol']}</div>", unsafe_allow_html=True)
        
        fig_hourly = px.bar(hourly_vol, x='hour_of_day', y='count', title="Messages by Hour of Day", template=plotly_template, color_discrete_sequence=custom_palette)
        fig_hourly.update_layout(xaxis=dict(tickmode='linear', tick0=0, dtick=1))
        c2.plotly_chart(fig_hourly, use_container_width=True)
        c2.markdown(f"<div class='insight-card'><b>Analysis Insight:</b> {insights['hourly_vol']}</div>", unsafe_allow_html=True)
        
        fig_heatmap = px.imshow(
            heatmap_pivot, 
            labels=dict(x="Hour of Day", y="Day of Week", color="Messages"),
            x=heatmap_pivot.columns,
            y=heatmap_pivot.index,
            color_continuous_scale="RdYlGn_r",
            template=plotly_template,
            title="Activity Heatmap (Day vs Hour)"
        )
        st.plotly_chart(fig_heatmap, use_container_width=True)
        st.markdown(f"<div class='insight-card'><b>Analysis Insight:</b> {insights['heatmap']}</div>", unsafe_allow_html=True)

    with tab2:
        c3, c4 = st.columns(2)
        fig_user_vol = px.bar(user_vol, x='count', y='user_name', orientation='h', title="Message Count per User", template=plotly_template, color_discrete_sequence=custom_palette)
        c3.plotly_chart(fig_user_vol, use_container_width=True)
        c3.markdown(f"<div class='insight-card'><b>Analysis Insight:</b> {insights['user_vol']}</div>", unsafe_allow_html=True)
        
        fig_share = px.pie(user_vol, values='count', names='user_name', title="Conversation Share (%)", template=plotly_template, hole=0.3, color_discrete_sequence=custom_palette)
        c4.plotly_chart(fig_share, use_container_width=True)
        c4.markdown(f"<div class='insight-card'><b>Analysis Insight:</b> {insights['user_share']}</div>", unsafe_allow_html=True)
        
        st.markdown("### User Ranking & Statistics")
        st.dataframe(user_stats, use_container_width=True)
        
    with tab3:
        c5, c6 = st.columns(2)
        fig_sent_time = px.line(daily_sentiment, x='full_date', y='rolling_7d', title="Rolling 7-Day Average Sentiment", template=plotly_template, color_discrete_sequence=custom_palette)
        c5.plotly_chart(fig_sent_time, use_container_width=True)
        c5.markdown(f"<div class='insight-card'><b>Analysis Insight:</b> {insights['sentiment_trend']}</div>", unsafe_allow_html=True)
        
        fig_sent_user = px.bar(user_sent, x='user_name', y='sentiment_score', title="Average Sentiment by User", template=plotly_template, color='sentiment_score', color_continuous_scale="RdBu")
        c6.plotly_chart(fig_sent_user, use_container_width=True)
        c6.markdown(f"<div class='insight-card'><b>Analysis Insight:</b> {insights['sentiment_compare']}</div>", unsafe_allow_html=True)
        
        c7, c8 = st.columns(2)
        fig_starters = px.pie(starters, values='count', names='user_name', title="Conversation Starters", template=plotly_template, hole=0.3, color_discrete_sequence=custom_palette)
        c7.plotly_chart(fig_starters, use_container_width=True)
        c7.markdown(f"<div class='insight-card'><b>Analysis Insight:</b> {insights['starters']}</div>", unsafe_allow_html=True)
        
        fig_weekend = px.bar(weekend_vol, x='day_type', y='count', title="Weekend vs Weekday", template=plotly_template, color='day_type', color_discrete_sequence=custom_palette)
        c8.plotly_chart(fig_weekend, use_container_width=True)
        c8.markdown(f"<div class='insight-card'><b>Analysis Insight:</b> {insights['weekend_vol']}</div>", unsafe_allow_html=True)
        
    with tab4:
        st.markdown("### Export BI Presentations & Reports")
        
        kpis = {
            "Total Messages": total_messages,
            "Unique Users": unique_users,
            "Average Messages / Day": avg_msg_per_day,
            "Most Active User": most_active_user,
            "Peak Activity Hour": peak_hour_str,
            "Average Sentiment Score": round(avg_sentiment, 2)
        }
        
        c_pdf, c_ppt = st.columns(2)
        
        with c_pdf:
            st.markdown("#### PDF Report")
            st.write("Download a detailed and fully labeled PDF document containing all tables, charts, and analysis conclusions.")
            if st.button("Generate PDF Report", type="primary", use_container_width=True):
                with st.spinner("Compiling PDF Report..."):
                    try:
                        pdf_bytes = generate_pdf_report_with_fallback(kpis, data_frames, user_stats, insights)
                        b64 = base64.b64encode(pdf_bytes).decode()
                        href = f'<a href="data:application/pdf;base64,{b64}" download="ChatPulse_Report.pdf" class="btn" style="display:block; text-align:center; padding:10px; background:#1B4EF5; color:white; border-radius:8px; text-decoration:none; font-weight:bold;">Download PDF Report</a>'
                        st.markdown(href, unsafe_allow_html=True)
                        st.success("PDF generated successfully!")
                    except Exception as e:
                        st.error(f"Failed to generate PDF: {str(e)}")
                        
        with c_ppt:
            st.markdown("#### PowerPoint Presentation")
            st.write("Generate a formatted slide deck matching the dashboard layout, visual themes, and insights.")
            if not pptx_installed:
                st.warning("The 'python-pptx' library is missing or could not be loaded. Please run 'pip install python-pptx' to activate presentation slide generation.")
            else:
                if st.button("Generate PPTX Slides", type="primary", use_container_width=True):
                    with st.spinner("Creating Presentation Slides..."):
                        try:
                            pptx_bytes = generate_pptx_report(kpis, data_frames, user_stats, insights)
                            b64 = base64.b64encode(pptx_bytes).decode()
                            href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="ChatPulse_Presentation.pptx" class="btn" style="display:block; text-align:center; padding:10px; background:#3874FF; color:white; border-radius:8px; text-decoration:none; font-weight:bold;">Download Presentation</a>'
                            st.markdown(href, unsafe_allow_html=True)
                            st.success("PowerPoint Presentation generated successfully!")
                        except Exception as e:
                            st.error(f"Failed to generate Presentation: {str(e)}")

else:
    if uploaded_file is None and not use_sample:
        st.info("Upload a WhatsApp export file (.txt or .zip) from the sidebar to begin, or click 'Use Sample Data Demo' to see it in action.")
